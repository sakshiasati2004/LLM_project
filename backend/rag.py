import os
from dotenv import load_dotenv
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    CSVLoader,
    UnstructuredWordDocumentLoader
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings
from langchain_core.documents import Document

load_dotenv()

FAISS_BASE_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "faiss_indexes"
)
os.makedirs(FAISS_BASE_DIR, exist_ok=True)


# -------------------- EXISTING LOADERS --------------------

def load_msg_file(file_path):
    """Load .msg (Outlook email) files"""
    try:
        import extract_msg
        from bs4 import BeautifulSoup

        msg = extract_msg.Message(file_path)

        body = msg.body
        if not body or not body.strip():
            if msg.htmlBody:
                html = msg.htmlBody
                if isinstance(html, bytes):
                    html = html.decode('utf-8', errors='ignore')
                soup = BeautifulSoup(html, "html.parser")
                body = soup.get_text(separator="\n").strip()

        subject = msg.subject.replace('\x00', '').strip() if msg.subject else ''

        content = f"""Subject: {subject}
From: {msg.sender}
To: {msg.to}
Date: {msg.date}
Body:
{body}
"""
        return [Document(page_content=content, metadata={"source": file_path})]
    except Exception as e:
        raise ValueError(f"Error loading .msg file: {str(e)}")


def load_chm_file(file_path):
    """Load .chm (Help) files"""
    try:
        import chm.chm as chmlib
        chm_file = chmlib.CHMFile()
        chm_file.LoadCHM(file_path)

        contents = []
        def get_files(chm, ui, ctx):
            path = ui.path.decode("utf-8") if isinstance(ui.path, bytes) else ui.path
            if path.endswith(".html") or path.endswith(".htm"):
                result, content = chm.RetrieveObject(ui)
                if result == 0 and content:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(content, "html.parser")
                    text = soup.get_text(separator="\n")
                    if text.strip():
                        contents.append(text)

        chm_file.EnumerateFiles(get_files, None)
        chm_file.CloseCHM()

        if not contents:
            raise ValueError("No readable content found in .chm file")

        full_text = "\n\n".join(contents)
        return [Document(page_content=full_text, metadata={"source": file_path})]
    except Exception as e:
        raise ValueError(f"Error loading .chm file: {str(e)}")


# -------------------- NEW LOADERS --------------------

def extract_text_from_image(image):
    """Extract text from a PIL image using Tesseract OCR"""
    try:
        import pytesseract
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        return ""


def load_image_file(file_path):
    """Load standalone image files (JPEG, PNG) using Tesseract OCR"""
    try:
        from PIL import Image
        import pytesseract

        image = Image.open(file_path)
        text = pytesseract.image_to_string(image).strip()

        if not text:
            text = "No readable text found in this image."

        return [Document(
            page_content=text,
            metadata={"source": file_path}
        )]
    except Exception as e:
        raise ValueError(f"Error loading image file: {str(e)}")


def load_pdf_with_images(file_path):
    """Load PDF — extract both text and images (OCR on embedded images)"""
    try:
        import fitz  # pymupdf
        from PIL import Image
        import io

        documents = []
        pdf = fitz.open(file_path)

        for page_num in range(len(pdf)):
            page = pdf[page_num]

            # Extract normal text
            text = page.get_text().strip()
            if text:
                documents.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "page": page_num}
                ))

            # Extract images from page and run OCR
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = pdf.extract_image(xref)
                image_bytes = base_image["image"]

                image = Image.open(io.BytesIO(image_bytes))
                ocr_text = extract_text_from_image(image)

                if ocr_text and len(ocr_text) > 10:
                    documents.append(Document(
                        page_content=ocr_text,
                        metadata={
                            "source": file_path,
                            "page": page_num,
                            "type": "image_ocr"
                        }
                    ))

        pdf.close()
        return documents if documents else PyPDFLoader(file_path).load()

    except Exception as e:
        return PyPDFLoader(file_path).load()


def load_docx_with_images(file_path):
    """Load DOCX — extract both text and images (OCR on embedded images)"""
    try:
        from docx import Document as DocxDocument
        from PIL import Image
        import io

        doc = DocxDocument(file_path)
        documents = []

        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if full_text:
            documents.append(Document(
                page_content=full_text,
                metadata={"source": file_path}
            ))

        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_part = rel.target_part
                image_bytes = image_part.blob
                try:
                    image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = extract_text_from_image(image)
                    if ocr_text and len(ocr_text) > 10:
                        documents.append(Document(
                            page_content=ocr_text,
                            metadata={
                                "source": file_path,
                                "type": "image_ocr"
                            }
                        ))
                except Exception:
                    continue

        return documents if documents else UnstructuredWordDocumentLoader(file_path).load()

    except Exception as e:
        return UnstructuredWordDocumentLoader(file_path).load()


def load_ppt_file(file_path):
    """Load PPT/PPTX — extract text from slides + OCR on slide images"""
    try:
        from pptx import Presentation
        from PIL import Image
        import io

        prs = Presentation(file_path)
        documents = []

        for slide_num, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)

                if shape.shape_type == 13:
                    try:
                        image_bytes = shape.image.blob
                        image = Image.open(io.BytesIO(image_bytes))
                        ocr_text = extract_text_from_image(image)
                        if ocr_text and len(ocr_text) > 10:
                            slide_text.append(f"[Image Text]: {ocr_text}")
                    except Exception:
                        continue

            if slide_text:
                documents.append(Document(
                    page_content="\n".join(slide_text),
                    metadata={
                        "source": file_path,
                        "slide": slide_num + 1
                    }
                ))

        return documents if documents else []

    except Exception as e:
        raise ValueError(f"Error loading PPT file: {str(e)}")


# -------------------- LOAD DOCUMENT --------------------

def load_document(file_path):
    """Load document based on file extension"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return load_pdf_with_images(file_path)
    elif ext == ".txt":
        return TextLoader(file_path).load()
    elif ext == ".csv":
        return CSVLoader(file_path).load()
    elif ext in [".doc", ".docx"]:
        return load_docx_with_images(file_path)
    elif ext == ".msg":
        return load_msg_file(file_path)
    elif ext == ".chm":
        return load_chm_file(file_path)
    elif ext in [".jpg", ".jpeg", ".png"]:
        return load_image_file(file_path)
    elif ext in [".ppt", ".pptx"]:
        return load_ppt_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def load_and_split(file_path):
    documents = load_document(file_path)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    chunks = splitter.split_documents(documents)

    valid_chunks = []
    for chunk in chunks:
        content = chunk.page_content.strip()
        if not content:
            continue
        if len(content) < 10:
            continue
        chunk.page_content = content.encode('utf-8', errors='ignore').decode('utf-8')
        valid_chunks.append(chunk)

    return valid_chunks


def add_metadata(chunks, user_id, session_id, file_name):
    for chunk in chunks:
        chunk.metadata.update({
            "user_id": str(user_id),
            "session_id": str(session_id),
            "file_name": file_name
        })
    return chunks


def get_embeddings():
    return OpenAIEmbeddings(
        api_key=os.getenv("OPENAI_API_KEY"),
        base_url="https://openrouter.ai/api/v1"
    )


def get_user_vectorstore_path(user_id):
    return os.path.join(FAISS_BASE_DIR, f"faiss_{user_id}")


def create_or_load_vectorstore(chunks, user_id):
    path = get_user_vectorstore_path(user_id)
    embeddings = get_embeddings()
    if os.path.exists(path):
        vectorstore = FAISS.load_local(
            path, embeddings, allow_dangerous_deserialization=True
        )
        vectorstore.add_documents(chunks)
    else:
        vectorstore = FAISS.from_documents(chunks, embeddings)
    vectorstore.save_local(path)
    return vectorstore


def load_existing_vectorstore(user_id):
    path = get_user_vectorstore_path(user_id)
    if not os.path.exists(path):
        return None
    embeddings = get_embeddings()
    return FAISS.load_local(
        path, embeddings, allow_dangerous_deserialization=True
    )


def get_context_from_query(vectorstore, query, user_id, session_id, selected_doc="All Documents"):
    """
    Returns: (context_string, docs_list, file_names_list)
    ✅ CHANGED: now returns actual docs list as 2nd element instead of "current"/"none" string
    This allows chat.py to inspect per-document sources for multi-doc answers
    """
    docs = vectorstore.similarity_search(query, k=10)

    if selected_doc != "All Documents":
        docs = [d for d in docs if d.metadata.get("file_name") == selected_doc]

    session_docs = [
        d for d in docs
        if str(d.metadata.get("session_id")) == str(session_id)
        and str(d.metadata.get("user_id")) == str(user_id)
    ]

    if session_docs:
        context = "\n\n".join([d.page_content for d in session_docs[:5]])
        file_names = list(set([
            d.metadata.get("file_name", "Unknown") for d in session_docs
        ]))
        # ✅ Return actual docs list (not "current" string) so chat.py can group by source
        return context, session_docs, file_names

    return "", [], []


def get_all_documents(vectorstore, user_id):
    if not vectorstore:
        return []
    try:
        docs = vectorstore.docstore._dict.values()
        return list(set([
            doc.metadata.get("file_name", "Unknown")
            for doc in docs
            if str(doc.metadata.get("user_id")) == str(user_id)
        ]))
    except Exception:
        return []


def get_session_documents(vectorstore, user_id, session_id):
    if not vectorstore:
        return []
    try:
        docs = vectorstore.docstore._dict.values()
        return list(set([
            doc.metadata.get("file_name", "Unknown")
            for doc in docs
            if str(doc.metadata.get("user_id")) == str(user_id)
            and str(doc.metadata.get("session_id")) == str(session_id)
        ]))
    except Exception:
        return []